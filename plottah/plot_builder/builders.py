from typing import Dict
from collections import defaultdict

import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt

from plottah.plot_handler import PlotHandler
from plottah.colors import PlotColors
from plottah.plot_builder.specific_builders import PLOT_BUILDERS_DICT


def build_univariate_plot(
    df,
    feature_col: str,
    target: str,
    feature_type: str = "float",
    specs: list = None,
    colors: PlotColors = PlotColors(),
    show_plot: bool = False,
    hoverinfo="all",
    n_bins: int = 10,
    bins: list = None,
    distplot_q_min: float = None,
    distplot_q_max: float = None,
):
    """
    buils standard univariate plot from days 'ye

    Returns
    """

    # if feature type not-categoriocal, assume numerical plot type
    feature_type = feature_type if feature_type == "categorical" else "numerical"

    # get appropriate plot builder
    plot_builder = PLOT_BUILDERS_DICT[feature_type]

    print(distplot_q_min)
    print(distplot_q_max)

    # build plot
    plot = plot_builder(
        df,
        feature_col=feature_col,
        target=target,
        feature_type=feature_type,
        colors=colors,
        show_plot=show_plot,
        hoverinfo=hoverinfo,
        n_bins=n_bins,
        bins=bins,
        specs=specs,
        distplot_q_min=distplot_q_min,
        distplot_q_max=distplot_q_max,
    )

    # show the plot if show_plot set to true
    if show_plot:
        plot.show()

    return plot


def build_univariate_plots(
    df,
    features: list,
    target: str,
    feature_types: dict,
    n_bins: dict = None,
    bins: dict = None,
    distplot_q_min: dict = None,
    distplot_q_max: dict = None,
    save_directory: pathlib.Path() = None,
    colors: PlotColors = PlotColors(),
    show_plot: bool = False,
    hoverinfo="none",
) -> Dict[str, PlotHandler]:
    """
    function that generates standard univariate plots

    Args:
        df (pd.DataFrame): dataframe containing columns to analyze,
        features (list): list of columns,
        save_directory (pathlib.Path()): where to store figures:

    Returns:
        Dict: map from feature name to figure

    """

    # print(distplot_q_max)
    # print(distplot_q_min)

    # if only a single feature passed wrap in a list
    if isinstance(features, str):
        features = [features]

    # create mapping from bins to none if not passed as argument
    if bins is None:
        bins = {feature: None for feature in features}

    # create mapping from features to 10 if nbins not passed as argument
    if n_bins is None:
        n_bins = {feature: 10 for feature in features}
    else:
        # n_bins = defaultdict(lambda: 10, n_bins)
        for feature in features:
            # user does not have to provide complete mapping
            if feature not in n_bins.keys():
                n_bins[feature] = 10

    # create mapping from features to float if type not provided
    if feature_types is None:
        feature_types = {feature: "float" for feature in features}

    # test - need to return None if user did not provive q min for all features but only some
    if distplot_q_min is None:
        distplot_q_min = {feature: None for feature in features}
    else:
        # user does not have to provide complete mapping
        distplot_q_min = defaultdict(lambda: None, distplot_q_min)

    # test - need to return None if user did not provive q min for all features but only some
    if distplot_q_max is None:
        distplot_q_max = {feature: None for feature in features}
    else:
        # user does not have to provide complete mapping
        distplot_q_max = defaultdict(lambda: None, distplot_q_max)

    # Run loop
    figs = {}
    save_locs = []
    for i, feature in enumerate(features):
        print(f"[{i+1}/{len(features)}] Starting univariate analysis for: {feature}")
        if feature not in df.columns:
            raise ValueError(f"{feature} not in columns of dataframe")

        fig = build_univariate_plot(
            df=df,
            feature_col=feature,
            target=target,
            feature_type=feature_types[feature],
            colors=colors,
            show_plot=show_plot,
            hoverinfo=hoverinfo,
            n_bins=n_bins[feature],
            bins=bins[feature],
            distplot_q_min=distplot_q_min[feature],
            distplot_q_max=distplot_q_max[feature],
        )

        if save_directory is not None:
            save_loc = fig.save_fig(save_directory)
            print(f"[{i+1}/{len(features)}] Saving univariate anaylsis for {feature}")
            save_locs.append(save_loc)

        figs[feature] = fig

    return figs, save_locs


def build_powerpoint(
    fig_locs: list,
    feature_names: list,
    save_path: pathlib.Path(),
):
    pres = Presentation()

    for i, fig_loc in enumerate(fig_locs):
        s_register = pres.slide_layouts[7]
        s = pres.slides.add_slide(s_register)

        pic = s.shapes.add_picture(
            str(pathlib.Path(fig_loc).resolve()),
            Inches(0.5),
            Inches(1.75),
            width=Inches(7),
            height=Inches(5),
        )
        pic.left = int((pres.slide_width - pic.width) / 2)

        title = s.shapes.title
        title.text = f"{i} Univariate analysis for {(feature_names)[i]}"
        title_para = s.shapes.title.text_frame.paragraphs[0]
        title_para.font.size = Pt(24)

    pres.save(save_path)
